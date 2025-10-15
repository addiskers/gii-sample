from flask import Flask, render_template, request, jsonify, send_file
import os
import json
import ast
import re
import traceback
import logging
import logging.handlers
import time
from datetime import datetime
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
from openai import OpenAI
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LEGEND_POSITION,
    XL_TICK_LABEL_POSITION,
    XL_TICK_MARK,
)
from enum import Enum
from typing import Optional, Dict, List, Any
import concurrent.futures
import threading

BALL_VALUES = {"High": 3, "Medium": 2, "Low": 1}
BALL_COLOR = RGBColor(108, 63, 197)  # Purple
EMPTY_COLOR = RGBColor(255, 255, 255)  # White


def setup_logging():
    """Setup logging configuration with UTF-8 encoding"""
    log_dir = "/app/logs" if os.path.exists("/app/logs") else "logs"
    os.makedirs(log_dir, exist_ok=True)

    app_log_file = os.path.join(log_dir, "app.log")
    timing_log_file = os.path.join(log_dir, "timing.log")

    logging.basicConfig(level=logging.INFO)

    main_logger = logging.getLogger("ppt_generator")
    main_logger.setLevel(logging.INFO)

    timing_logger = logging.getLogger("timing")
    timing_logger.setLevel(logging.INFO)

    if not main_logger.handlers:
        formatter = logging.Formatter(
            "%(asctime)s - %(name)s - %(levelname)s - %(message)s",
            datefmt="%Y-%m-%d %H:%M:%S",
        )

        app_handler = logging.handlers.RotatingFileHandler(
            app_log_file, maxBytes=10 * 1024 * 1024, backupCount=5, encoding='utf-8'
        )
        app_handler.setFormatter(formatter)
        main_logger.addHandler(app_handler)

        console_handler = logging.StreamHandler()
        console_handler.setFormatter(formatter)
        if hasattr(console_handler.stream, 'reconfigure'):
            console_handler.stream.reconfigure(encoding='utf-8')
        main_logger.addHandler(console_handler)

    if not timing_logger.handlers:
        timing_formatter = logging.Formatter(
            "%(asctime)s - %(message)s", datefmt="%Y-%m-%d %H:%M:%S"
        )

        timing_handler = logging.handlers.RotatingFileHandler(
            timing_log_file, maxBytes=5 * 1024 * 1024, backupCount=3, encoding='utf-8'
        )
        timing_handler.setFormatter(timing_formatter)
        timing_logger.addHandler(timing_handler)

    return main_logger, timing_logger


logger, timing_logger = setup_logging()

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 20 * 1024 * 1024
app.config["UPLOAD_FOLDER"] = "generated_ppts"
app.config["TEMPLATE_FOLDER"] = "templates"

os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
os.makedirs(app.config["TEMPLATE_FOLDER"], exist_ok=True)

logger.info("Application starting up")
logger.info(f"Upload folder: {app.config['UPLOAD_FOLDER']}")
logger.info(f"Template folder: {app.config['TEMPLATE_FOLDER']}")

load_dotenv()
api_key_openAI = os.getenv("OPENAI_API_KEY")
api_key_gemini = os.getenv("GEMINI_API_KEY")

if api_key_openAI:
    logger.info("OpenAI API key loaded successfully")
else:
    logger.warning("OpenAI API key not found")

if api_key_gemini:
    logger.info("Gemini API key loaded successfully")
else:
    logger.warning("Gemini API key not found")

client = OpenAI(api_key=api_key_openAI)
genai.configure(api_key=api_key_gemini)

# ==================== SUCCESS FACTORS FUNCTIONS ====================


def replace_text_preserve_format(paragraph, old_text, new_text):
    """Replace text in a paragraph while preserving formatting"""
    full_text = paragraph.text
    
    if old_text not in full_text:
        return False
    
    runs = list(paragraph.runs)
    
    for run in runs:
        if old_text in run.text:
            run.text = run.text.replace(old_text, new_text)
            return True
    
    new_full_text = full_text.replace(old_text, new_text)
    
    if runs:
        first_run = runs[0]
        font_info = {
            'name': first_run.font.name,
            'size': first_run.font.size,
            'bold': first_run.font.bold,
            'italic': first_run.font.italic,
            'color': first_run.font.color.rgb if first_run.font.color.type == 1 else None
        }
        
        for run in runs:
            run.text = ''
        
        first_run.text = new_full_text
        
        if font_info['name']:
            first_run.font.name = font_info['name']
        if font_info['size']:
            first_run.font.size = font_info['size']
        if font_info['bold'] is not None:
            first_run.font.bold = font_info['bold']
        if font_info['italic'] is not None:
            first_run.font.italic = font_info['italic']
        if font_info['color']:
            first_run.font.color.rgb = font_info['color']
    else:
        paragraph.text = new_full_text
    
    return True


def apply_success_factors_to_slide(slide, success_factors_data):
    """Apply pre-generated success factors data to slide"""
    logger.info("Applying success factors to slide")
    
    market_name = success_factors_data["market_name"]
    titles = [success_factors_data[f"ksfactor{i+1}_title"] for i in range(7)]
    examples = [success_factors_data[f"ksfactor{i+1}_example"] for i in range(7)]
    titles_list = success_factors_data["titles"]
    
    replacements = 0
    
    for shape_idx, shape in enumerate(slide.shapes):
        
        if shape.has_text_frame:
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                full_text = paragraph.text
                
                if "{{market_name}}" in full_text:
                    if replace_text_preserve_format(paragraph, "{{market_name}}", market_name):
                        replacements += 1
                
                if "{{titles}}" in full_text:
                    if replace_text_preserve_format(paragraph, "{{titles}}", titles_list):
                        replacements += 1
                
                for i in range(7):
                    title_placeholder = f"{{{{ksfactor{i+1}_title}}}}"
                    if title_placeholder in full_text:
                        if replace_text_preserve_format(paragraph, title_placeholder, titles[i]):
                            replacements += 1
                
                for i in range(7):
                    example_placeholder = f"{{{{ksfactor{i+1}_example}}}}"
                    if example_placeholder in full_text:
                        if replace_text_preserve_format(paragraph, example_placeholder, examples[i]):
                            replacements += 1
        
        # Check tables
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            full_text = paragraph.text
                            
                            if "{{market_name}}" in full_text:
                                replace_text_preserve_format(paragraph, "{{market_name}}", market_name)
                                replacements += 1
                            
                            if "{{titles}}" in full_text:
                                replace_text_preserve_format(paragraph, "{{titles}}", titles_list)
                                replacements += 1
                            
                            for i in range(7):
                                title_placeholder = f"{{{{ksfactor{i+1}_title}}}}"
                                example_placeholder = f"{{{{ksfactor{i+1}_example}}}}"
                                
                                if title_placeholder in full_text:
                                    replace_text_preserve_format(paragraph, title_placeholder, titles[i])
                                    replacements += 1
                                
                                if example_placeholder in full_text:
                                    replace_text_preserve_format(paragraph, example_placeholder, examples[i])
                                    replacements += 1
        
        # Check charts
        if shape.has_chart:
            try:
                chart = shape.chart
                for series in chart.series:
                    for point in series.points:
                        try:
                            if hasattr(point, 'data_label') and point.data_label:
                                label = point.data_label
                                if hasattr(label, 'text_frame') and label.text_frame:
                                    for para_idx, paragraph in enumerate(label.text_frame.paragraphs):
                                        full_text = paragraph.text
                                        
                                        if "{{market_name}}" in full_text:
                                            replace_text_preserve_format(paragraph, "{{market_name}}", market_name)
                                            replacements += 1
                                        
                                        if "{{titles}}" in full_text:
                                            replace_text_preserve_format(paragraph, "{{titles}}", titles_list)
                                            replacements += 1
                                        
                                        for j in range(7):
                                            title_placeholder = f"{{{{ksfactor{j+1}_title}}}}"
                                            example_placeholder = f"{{{{ksfactor{j+1}_example}}}}"
                                            
                                            if title_placeholder in full_text:
                                                replace_text_preserve_format(paragraph, title_placeholder, titles[j])
                                                replacements += 1
                                            
                                            if example_placeholder in full_text:
                                                replace_text_preserve_format(paragraph, example_placeholder, examples[j])
                                                replacements += 1
                                        
                                        try:
                                            for run in paragraph.runs:
                                                if para_idx == 0:
                                                    run.font.bold = True
                                                elif para_idx == 1:
                                                    run.font.bold = False
                                        except:
                                            pass
                        except:
                            pass
            except:
                pass
    
    logger.info(f"✓ Total success factors replacements made: {replacements}")


# ==================== END SUCCESS FACTORS FUNCTIONS ====================

# ==================== DEFINITION TABLE FUNCTIONS ====================


def build_definition_batch_prompt(nested_dict: dict, headline: str) -> str:
    """Build a single batch prompt for all definitions - OPTIMIZED"""
    headline_clean = headline.replace("Global ", "").replace("global ", "").replace("Market ", "").replace("market ", "")
    definitions_needed = []
    
    for level_0_name, level_1_dict in nested_dict.items():
        all_segments_in_category = list(level_1_dict.keys())
        
        for level_1_name, level_2_dict in level_1_dict.items():
            level_2_names = list(level_2_dict.keys()) if level_2_dict else []
            sub_areas = (
                ", ".join(level_2_names) if level_2_names else "general concepts"
            )
            
            is_others = level_1_name.lower() in ["others", "other"]
            
            if is_others:
                sibling_segments = [s for s in all_segments_in_category if s.lower() not in ["others", "other"]]
                definitions_needed.append({
                    "name": level_1_name,
                    "sub_areas": sub_areas,
                    "is_others": True,
                    "parent_category": level_0_name,
                    "sibling_segments": sibling_segments
                })
            else:
                definitions_needed.append({
                    "name": level_1_name,
                    "sub_areas": sub_areas,
                    "is_others": False
                })

    prompt = f"""Generate brief definitions for {headline}.

Overall definition:
{headline_clean}: [Write 1-2 concise sentences]

Segment definitions (1-2 sentences each):
"""

    for item in definitions_needed:
        if item.get("is_others", False):
            prompt += f"\n{item['name']} (in {item['parent_category']} category): [Define what other {item['parent_category'].lower()} are included beyond {', '.join(item['sibling_segments'])}. Be specific about what additional categories are covered.]"
        else:
            prompt += f"\n{item['name']}: [Brief definition based on {item['sub_areas']}]"

    prompt += f"""

Return ONLY valid JSON (no extra text):
{{
    "{headline_clean}": "definition here",
    "Segment Name": "definition here"
}}

Requirements:
- Each definition: 1-2 sentences ONLY
- Professional, concise language
- For "Others" segments: Specifically explain what additional items/categories are included beyond the explicitly listed segments in that category
- No extra newlines or spacing
"""

    return prompt

def generate_all_definitions_batch(nested_dict: dict, headline: str) -> dict:
    """Generate all definitions in ONE API call - OPTIMIZED"""
    logger.info("Generating all definitions in a single API call...")

    prompt = build_definition_batch_prompt(nested_dict, headline)

    response = client.chat.completions.create(
        model="gpt-5-mini",
        messages=[
            {
                "role": "system",
                "content": "You are a market research expert. Generate concise, professional definitions in JSON format. Keep each definition to 1-2 sentences maximum.",
            },
            {"role": "user", "content": prompt},
        ],
        response_format={"type": "json_object"},
    )

    definitions = json.loads(response.choices[0].message.content)
    
    cleaned_definitions = {}
    for key, value in definitions.items():
        cleaned_value = ' '.join(value.split())
        cleaned_definitions[key] = cleaned_value
    
    logger.info(f"Generated {len(cleaned_definitions)} definitions in one call")
    return cleaned_definitions


def build_table_data(nested_dict: dict, definitions: dict, headline: str) -> list:
    """Build table sections from definitions"""
    table_sections = []

    headline_clean = headline.replace("Global ", "").replace("global ", "").replace("Market ", "").replace("market ", "")
    headline_def = definitions.get(headline_clean, "Market definition")

    table_sections.append(
        {
            "header": "DEFINITION",
            "rows": [{"title": headline_clean, "definition": headline_def}],
        }
    )

    for level_0_name, level_1_dict in nested_dict.items():
        table_rows = []
        for level_1_name in level_1_dict.keys():
            definition = definitions.get(level_1_name, f"Definition for {level_1_name}")
            table_rows.append({"title": level_1_name, "definition": definition})

        table_sections.append({"header": level_0_name.upper(), "rows": table_rows})

    return table_sections


def calculate_text_height(text, font_size_pt, cell_width_inches):
    """Calculate required cell height for text"""
    if not text or not text.strip():
        return Inches(0.3)

    chars_per_line = max(30, int((cell_width_inches * 72 - 10) / (font_size_pt * 0.55)))
    words = text.split()
    lines = []
    current_line = ""

    for word in words:
        test_line = f"{current_line} {word}" if current_line else word
        if len(test_line) <= chars_per_line:
            current_line = test_line
        else:
            if current_line:
                lines.append(current_line)
            current_line = word

    if current_line:
        lines.append(current_line)

    lines_needed = max(1, len(lines))
    base_height = (lines_needed * font_size_pt + 6) / 72
    return Inches(min(base_height, 2.5))


def create_table_on_slide(
    slide,
    table_section,
    table_top,
    table_width,
    left_pos,
    rows_subset,
    is_continuation,
    is_definition=False,
):
    """Create table on slide"""
    header_height = Inches(0.4)
    row_heights = [header_height]
    right_col_width = (table_width * 0.75) / 914400

    for row in rows_subset:
        calculated_height = calculate_text_height(row["definition"], 9, right_col_width)
        row_heights.append(calculated_height + Inches(0.20))

    shape = slide.shapes.add_table(
        len(rows_subset) + 1, 2, left_pos, table_top, table_width, sum(row_heights)
    )
    table = shape.table
    table.columns[0].width = int(table_width * 0.25)
    table.columns[1].width = int(table_width * 0.75)

    table.first_row = False
    table.first_col = False

    for i, height in enumerate(row_heights):
        table.rows[i].height = int(height)

    # Header cell
    header_cell = table.cell(0, 0)
    header_cell.merge(table.cell(0, 1))
    header_cell.text = (
        f"{table_section['header']}{' (continued)' if is_continuation else ''}"
    )
    header_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    header_cell.fill.solid()

    if is_definition:
        header_cell.fill.fore_color.rgb = RGBColor(49, 9, 126)
    else:
        header_cell.fill.fore_color.rgb = RGBColor(251, 89, 45)

    frame = header_cell.text_frame
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = (
        Inches(0.05)
    )
    para = frame.paragraphs[0]
    para.font.size = Pt(12)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)
    para.alignment = PP_ALIGN.CENTER

    # Content rows
    for row_idx, row_data in enumerate(rows_subset, 1):
        left_cell = table.cell(row_idx, 0)
        left_cell.text = row_data["title"]
        left_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        left_cell.fill.solid()
        left_cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

        left_frame = left_cell.text_frame
        left_frame.margin_left = left_frame.margin_right = Inches(0.05)
        left_frame.margin_top = left_frame.margin_bottom = Inches(0.10)
        left_para = left_frame.paragraphs[0]
        left_para.font.size = Pt(9)
        left_para.font.bold = True
        left_para.alignment = PP_ALIGN.CENTER

        right_cell = table.cell(row_idx, 1)
        right_cell.text = row_data["definition"]
        right_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        right_cell.fill.solid()
        right_cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

        right_frame = right_cell.text_frame
        right_frame.margin_left = right_frame.margin_right = Inches(0.05)
        right_frame.margin_top = right_frame.margin_bottom = Inches(0.10)
        right_para = right_frame.paragraphs[0]
        right_para.font.size = Pt(9)
        right_para.alignment = PP_ALIGN.JUSTIFY

    return sum(row_heights)


def duplicate_slide(prs, slide_index, insert_after_index=None):
    """Duplicate slide and insert at specific position"""
    original_slide = prs.slides[slide_index]

    blank_slide = prs.slides.add_slide(original_slide.slide_layout)

    for shape in original_slide.shapes:
        if shape.has_text_frame and not shape.has_table:
            if shape.is_placeholder:
                try:
                    placeholder = blank_slide.placeholders[shape.placeholder_format.idx]
                    if shape.text.strip() and "click to add" not in shape.text.lower():
                        placeholder.text = shape.text
                except:
                    pass

    if insert_after_index is not None:
        slides = prs.slides._sldIdLst
        slide_to_move = slides[-1]
        slides.remove(slide_to_move)
        slides.insert(insert_after_index + 1, slide_to_move)

    return blank_slide


def insert_definition_tables(prs, nested_dict, headline, start_index=10):
    """Insert definition table slides starting at start_index"""
    logger.info(f"Inserting definition tables starting at slide {start_index}")

    start_time = time.time()
    definitions = generate_all_definitions_batch(nested_dict, headline)
    elapsed = time.time() - start_time
    timing_logger.info(f"Definition generation completed in {elapsed:.2f}s")

    table_data = build_table_data(nested_dict, definitions, headline)
    logger.info(f"Built {len(table_data)} table sections")

    margin = Inches(0.5)
    table_width = prs.slide_width - (2 * margin)
    top_start = Inches(1.1)
    table_spacing = Inches(0.1)
    max_slide_height = Inches(6.7)

    current_slide = prs.slides[start_index]
    current_top = top_start
    slides_created = 0
    current_insert_position = start_index

    for idx, section in enumerate(table_data):
        rows_remaining = section["rows"].copy()
        is_first_part = True
        is_definition = idx == 0 and section["header"] == "DEFINITION"
        right_col_width = (table_width * 0.75) / 914400

        while rows_remaining:
            available_space = max_slide_height - current_top
            header_height = Inches(0.35)
            first_row_height = calculate_text_height(
                rows_remaining[0]["definition"], 9, right_col_width
            ) + Inches(0.20)
            min_needed = header_height + first_row_height

            if available_space < min_needed:
                current_slide = duplicate_slide(
                    prs, start_index, current_insert_position
                )
                slides_created += 1
                current_insert_position += 1
                current_top = top_start
                available_space = max_slide_height - current_top

            rows_to_fit = []
            test_height = header_height

            for row in rows_remaining:
                row_height = calculate_text_height(
                    row["definition"], 9, right_col_width
                ) + Inches(0.20)
                if test_height + row_height <= available_space:
                    rows_to_fit.append(row)
                    test_height += row_height
                else:
                    break

            if not rows_to_fit:
                rows_to_fit = [rows_remaining[0]]

            temp_section = {"header": section["header"], "rows": rows_to_fit}
            table_height = create_table_on_slide(
                current_slide,
                temp_section,
                current_top,
                table_width,
                margin,
                rows_to_fit,
                not is_first_part,
                is_definition,
            )

            current_top += table_height + table_spacing
            rows_remaining = rows_remaining[len(rows_to_fit) :]
            is_first_part = False

    total_definition_slides = slides_created + 1
    logger.info(
        f"Created {total_definition_slides} definition slides (1 template + {slides_created} additional)"
    )

    return total_definition_slides


# ==================== END DEFINITION TABLE FUNCTIONS ====================

# ==================== MARKET IMPACTING FACTORS FUNCTIONS ====================


def sanitize_unicode_text(text):
    """Sanitize Unicode text by replacing special characters with ASCII equivalents"""
    if not text:
        return text
    
    replacements = {
        '\u2011': '-',
        '\u2013': '-',
        '\u2014': '--',
        '\u2018': "'",
        '\u2019': "'",
        '\u201c': '"',
        '\u201d': '"',
        '\u2022': '*',
        '\u2026': '...',
    }
    
    for unicode_char, ascii_char in replacements.items():
        text = text.replace(unicode_char, ascii_char)
    
    return text


def add_balls(slide, ball_positions, ball_values):
    """Add balls to specified positions on slide"""
    logger.info(f"Adding {len(ball_positions)} ball indicators to slide")
    
    for i, (left, top, width, height) in enumerate(ball_positions):
        if i < len(ball_values):
            value = ball_values[i]
            num_filled = BALL_VALUES.get(value, 0)
            num_total = 3
            
            ball_diameter = min(height * 0.6, width / 4)
            spacing = (width - num_total * ball_diameter) / (num_total + 1)
            
            for j in range(num_total):
                ball_left = left + spacing * (j + 1) + ball_diameter * j
                ball_top = top + (height - ball_diameter) / 2
                
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, ball_left, ball_top, ball_diameter, ball_diameter
                )
                
                if j < num_filled:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = BALL_COLOR
                else:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = EMPTY_COLOR
                
                shape.line.color.rgb = BALL_COLOR
                shape.line.width = Pt(2)
    
    logger.info("Ball indicators added successfully")


def parse_market_factors_response(response_text):
    """Parse AI response to extract headline, factors, and impact data"""
    if not response_text:
        return None, [], []
    
    lines = response_text.split('\n')
    headline = None
    factors = []
    table_data = []
    
    for line in lines:
        if 'Market Impacting Factors' in line:
            headline = sanitize_unicode_text(line.strip())
            break
    
    for line in lines:
        if line.strip().startswith(('Factor 1:', 'Factor 2:', 'Factor 3:', 'Factor 4:', 'Factor 5:')):
            factor = line.split(':', 1)[1].strip() if ':' in line else line.strip()
            factor = factor.replace('*', '').strip()
            factor = sanitize_unicode_text(factor)
            factors.append(factor)
    
    table_started = False
    for line in lines:
        if '|' in line and ('Factor' in line or table_started):
            if 'Factor' in line and 'North America' in line:
                table_started = True
                continue
            if '---' in line:
                continue
            if table_started and line.strip():
                parts = [p.strip() for p in line.split('|') if p.strip()]
                if len(parts) >= 6:
                    table_data.append(parts[1:6])
    
    return headline, factors, table_data


def apply_market_factors_to_slide(prs, slide_index, market_factors_data):
    """Apply pre-generated market factors data to slide"""
    logger.info(f"Applying market impacting factors to slide {slide_index}")
    
    headline = market_factors_data["headline"]
    factors = market_factors_data["factors"]
    impact_data = market_factors_data["impact_data"]
    
    slide = prs.slides[slide_index]
    
    # Replace headline
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if "{{headline}}" in paragraph.text:
                    original_runs = []
                    for run in paragraph.runs:
                        font_color = None
                        try:
                            if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                font_color = run.font.color.rgb
                        except (AttributeError, Exception):
                            pass
                        
                        original_runs.append({
                            'font_name': run.font.name,
                            'font_size': run.font.size,
                            'font_bold': run.font.bold,
                            'font_color': font_color
                        })
                    
                    paragraph.text = paragraph.text.replace("{{headline}}", headline)
                    
                    if original_runs and paragraph.runs:
                        for i, run in enumerate(paragraph.runs):
                            if i < len(original_runs):
                                if original_runs[i]['font_name']:
                                    run.font.name = original_runs[i]['font_name']
                                if original_runs[i]['font_size']:
                                    run.font.size = original_runs[i]['font_size']
                                if original_runs[i]['font_bold'] is not None:
                                    run.font.bold = original_runs[i]['font_bold']
                                if original_runs[i]['font_color']:
                                    try:
                                        run.font.color.rgb = original_runs[i]['font_color']
                                    except:
                                        pass
    
    # Replace factors in table
    replacements = 0
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text
                    
                    for i in range(5):
                        placeholder = f"{{{{factor{i+1}}}}}"
                        if placeholder in cell_text and i < len(factors):
                            cell.text = cell_text.replace(placeholder, factors[i])
                            
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(12)
                            
                            replacements += 1
                            cell_text = cell.text
    
    logger.info(f"Total factor replacements: {replacements}")
    
    # Add balls
    table_shape = None
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
            table = shape.table
            break
    
    if table_shape and impact_data:
        ball_positions = []
        ball_values = []
        
        table_left = table_shape.left
        table_top = table_shape.top
        
        for r in range(min(5, len(impact_data))):
            row_data = impact_data[r]
            data_row_index = r + 2
            
            if data_row_index < len(table.rows):
                for c in range(min(5, len(row_data))):
                    data_col_index = c + 1
                    
                    if data_col_index < len(table.columns):
                        cell_left = table_left
                        for col_idx in range(data_col_index):
                            cell_left += table.columns[col_idx].width
                        
                        cell_top = table_top
                        for row_idx in range(data_row_index):
                            cell_top += table.rows[row_idx].height
                        
                        cell_width = table.columns[data_col_index].width
                        cell_height = table.rows[data_row_index].height
                        
                        ball_positions.append((cell_left, cell_top, cell_width, cell_height))
                        ball_values.append(row_data[c])
        
        add_balls(slide, ball_positions, ball_values)
    
    logger.info(f"Market factors slide {slide_index} processed successfully")


# ==================== END MARKET IMPACTING FACTORS FUNCTIONS ====================


class AIRequestType(Enum):
    EXECUTIVE_SUMMARY = "executive_summary"
    MARKET_ENABLERS = "market_enablers"
    INDUSTRY_EXPANSION = "industry_expansion"
    INDUSTRY_EXPANSION_1 = "industry_expansion_1"
    INVESTMENT_CHALLENGES = "investment_challenges"
    COMPANY_INFO = "company_info"
    RESEARCH_JOURNALS = "research_journals"
    INDUSTRY_ASSOCIATIONS = "industry_associations"
    FINANCIAL_OVERVIEW = "financial_overview"
    SUCCESS_FACTORS = "success_factors"
    MARKET_IMPACTING_FACTORS = "market_impacting_factors"


class AIService:
    def __init__(self, openai_client, gemini_api_key):
        logger.info("Initializing AI Service")
        self.openai_client = openai_client
        self.gemini_configured = False
        if gemini_api_key:
            genai.configure(api_key=gemini_api_key)
            self.gemini_configured = True
            logger.info("Gemini configured successfully")
        else:
            logger.warning("Gemini not configured")

    def generate_content(
        self,
        request_type: AIRequestType,
        context: Dict[str, Any],
        existing_title: str = None,
    ) -> Any:
        start_time = time.time()
        logger.info(f"Generating content for: {request_type.value}")

        try:
            if request_type == AIRequestType.EXECUTIVE_SUMMARY:
                result = self._generate_executive_summary(context)
            elif request_type == AIRequestType.MARKET_ENABLERS:
                result = self._generate_market_enablers(context)
            elif request_type == AIRequestType.INDUSTRY_EXPANSION:
                result = self._generate_industry_expansion(context)
            elif request_type == AIRequestType.INDUSTRY_EXPANSION_1:
                result = self._generate_industry_expansion_1(context, existing_title)
            elif request_type == AIRequestType.INVESTMENT_CHALLENGES:
                result = self._generate_investment_challenges(context)
            elif request_type == AIRequestType.COMPANY_INFO:
                result = self._generate_company_info(context)
            elif request_type == AIRequestType.RESEARCH_JOURNALS:
                result = self._generate_research_journals(context)
            elif request_type == AIRequestType.INDUSTRY_ASSOCIATIONS:
                result = self._generate_industry_associations(context)
            elif request_type == AIRequestType.FINANCIAL_OVERVIEW:
                result = self._generate_financial_overview(context)
            elif request_type == AIRequestType.SUCCESS_FACTORS:
                result = self._generate_success_factors(context)
            elif request_type == AIRequestType.MARKET_IMPACTING_FACTORS:
                result = self._generate_market_impacting_factors(context)

            elapsed = time.time() - start_time
            timing_logger.info(f"{request_type.value} completed in {elapsed:.2f}s")
            logger.info(f"Content generation completed for: {request_type.value}")
            return result

        except Exception as e:
            elapsed = time.time() - start_time
            logger.error(
                f"Error generating {request_type.value} after {elapsed:.2f}s: {str(e)}"
            )
            raise

    def generate_content_parallel(self, ai_context: Dict[str, Any], nested_dict: Dict, headline: str) -> Dict[str, Any]:
        """OPTIMIZED: Run ALL tasks including definitions in parallel"""
        start_time = time.time()
        logger.info("Starting OPTIMIZED parallel AI content generation")
        results = {}

        phase1_tasks = {
            "executive_summary": (AIRequestType.EXECUTIVE_SUMMARY, ai_context, None, None),
            "market_enablers": (AIRequestType.MARKET_ENABLERS, ai_context, None, None),
            "industry_expansion": (AIRequestType.INDUSTRY_EXPANSION, ai_context, None, None),
            "investment_challenges": (AIRequestType.INVESTMENT_CHALLENGES, ai_context, None, None),
            "research_journals": (AIRequestType.RESEARCH_JOURNALS, ai_context, None, None),
            "industry_associations": (AIRequestType.INDUSTRY_ASSOCIATIONS, ai_context, None, None),
            "company_info": (AIRequestType.COMPANY_INFO, ai_context, None, None),
            "financial_overview": (AIRequestType.FINANCIAL_OVERVIEW, ai_context, None, None),
            "success_factors": (AIRequestType.SUCCESS_FACTORS, ai_context, None, None),
            "market_impacting_factors": (AIRequestType.MARKET_IMPACTING_FACTORS, ai_context, None, None),
            "definitions": (None, None, nested_dict, headline),
        }

        logger.info(f"Phase 1: Executing {len(phase1_tasks)} tasks in parallel")

        def execute_task(task_data):
            request_type, context, nested_d, hdline = task_data
            if request_type:  
                return self.generate_content(request_type, context)
            else: 
                return generate_all_definitions_batch(nested_d, hdline)

        with concurrent.futures.ThreadPoolExecutor(max_workers=20) as executor: 
            future_to_key = {
                executor.submit(execute_task, task_data): key
                for key, task_data in phase1_tasks.items()
            }

            for future in concurrent.futures.as_completed(future_to_key):
                key = future_to_key[future]
                try:
                    results[key] = future.result()
                    logger.info(f"Phase 1 task completed: {key}")
                except Exception as exc:
                    logger.error(f"Phase 1 task {key} generated an exception: {exc}")
                    raise exc

        logger.info("Phase 1 completed, starting Phase 2")
        
        industry_title = results["industry_expansion"]["title"]
        results["industry_expansion_1"] = self.generate_content(
            AIRequestType.INDUSTRY_EXPANSION_1, ai_context, industry_title
        )

        elapsed = time.time() - start_time
        timing_logger.info(f"OPTIMIZED parallel AI generation completed in {elapsed:.2f}s")
        logger.info("Parallel AI content generation completed successfully")
        return results

    def _generate_executive_summary(self, context: Dict[str, Any]) -> str:
        logger.info("Generating executive summary")
        first_line = (
            f"The {context['headline']} is valued at {context['cur']} {context['rev_current']} "
            f"{context['value_in']} in {context['base_year']}, and is expected to reach "
            f"{context['cur']} {context['rev_future']} {context['value_in']} by {context['forecast_year']}. "
            f"The market shows a steady CAGR of {context.get('cagr')}% from 2025 to 2032."
        )

        prompt = f"Write an executive summary for {context['headline']} focusing on key market drivers, trends, and growth factors within 50 words stricly. Do not include market size or revenue figures as they are already provided. Focus on qualitative insights about market dynamics, key players, and future outlook. ( start directly from setence without any intro like 'The executive summary is...')"

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini", messages=[{"role": "user", "content": prompt}]
        )
        ai_summary = response.choices[0].message.content
        full_summary = f"{first_line} {ai_summary}"
        logger.info("Executive summary generated successfully")
        return full_summary

    def _generate_market_enablers(self, context: Dict[str, Any]) -> str:
        logger.info("Generating market enablers")
        prompt = f'Write an executive summary about key market enablers (2 points) for {context["headline"]}, each 50 words strickly. Return a Python list like ["heading: context", "heading: context"].'
        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini", messages=[{"role": "user", "content": prompt}]
        )
        result = "\n".join(ast.literal_eval(response.choices[0].message.content))
        logger.info("Market enablers generated successfully")
        return result

    def _generate_industry_expansion(self, context: Dict[str, Any]) -> Dict[str, Any]:
        logger.info("Generating industry expansion")
        prompt = (
            f'Write one TOP Key Driver for the {context["headline"]} market. '
            f"Include a clear heading for the driver. "
            f"Return the output strictly as a Python dictionary with the following structure: "
            f'{{"title": "7–10 words", "paragraphs": ["paragraph1", "paragraph2", "paragraph3","parapragh4"]}}. '
            f"Each paragraph should be 80  words strict, qualitative in tone give 4 paragraphs, "
            f"and include real-world examples and facts. "
            f"Do not include market size, numbers, or links."
        )

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini", messages=[{"role": "user", "content": prompt}]
        )

        result = ast.literal_eval(response.choices[0].message.content)
        logger.info("Industry expansion generated successfully")
        return result

    def _generate_industry_expansion_1(
        self, context: Dict[str, Any], existing_title: str = None
    ) -> Dict[str, Any]:
        logger.info(
            f"Generating industry expansion 1 (avoiding title: {existing_title})"
        )
        existing_title_instruction = ""
        if existing_title:
            existing_title_instruction = f' Do not use "{existing_title}" as the title - generate a completely different driver.'

        prompt = (
            f'Write one TOP Key Driver for the {context["headline"]} market that is DIFFERENT from previous drivers.{existing_title_instruction} '
            f"Include a clear heading for the driver. "
            f"Return the output strictly as a Python dictionary with the following structure: "
            f'{{"title": "7–10 words", "paragraphs": ["paragraph1", "paragraph2", "paragraph3","parapragh4"]}}. '
            f"Each paragraph should be 80 words strict, qualitative in tone ,give 4 paragraphs must, "
            f"and include real-world examples and facts. "
            f"Do not include market size, numbers, or links. "
            f"Focus on a unique aspect not covered by other drivers."
        )

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini", messages=[{"role": "user", "content": prompt}]
        )

        result = ast.literal_eval(response.choices[0].message.content)
        logger.info("Industry expansion 1 generated successfully")
        return result

    def _generate_investment_challenges(self, context: Dict[str, Any]) -> str:
        logger.info("Generating investment challenges")
        prompt = (
            f'Write one TOP Key MARKET RESTRAINTS or CHALLENGES for the {context["headline"]} market. '
            f"Include a clear heading for the driver. "
            f"Return the output strictly as a Python dictionary with the following structure: "
            f'{{"title": "7–10 words", "paragraphs": ["paragraph1", "paragraph2", "paragraph3","parapragh4"]}}. '
            f"Each paragraph should be 80 words strict, qualitative in tone, give 4 paragraphs must"
            f"and include real-world examples and facts. "
            f"Do not include market size, numbers, or links."
        )

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini", messages=[{"role": "user", "content": prompt}]
        )

        result = ast.literal_eval(response.choices[0].message.content)
        logger.info("Investment challenges generated successfully")
        return result

    def _generate_success_factors(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Generate success factors data"""
        logger.info(f"Generating success factors for: {context['headline']}")
        
        prompt = f"""Generate exactly 7 key success factors for the {context['headline']}.

For each factor provide:
1. A short title (3-4 words max)
2. A very brief example (8-10 words maximum)

Format your response exactly like this:

Factor 1: [Title]
Example: [Brief Example]

Factor 2: [Title]
Example: [Brief Example]

Factor 3: [Title]
Example: [Brief Example]

Factor 4: [Title]
Example: [Brief Example]

Factor 5: [Title]
Example: [Brief Example]

Factor 6: [Title]
Example: [Brief Example]

Factor 7: [Title]
Example: [Brief Example]"""

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a market research analyst. Provide concise success factors in the specified format.",
                },
                {"role": "user", "content": prompt},
            ],
        )
        
        ai_response = response.choices[0].message.content.strip()
        
        lines = ai_response.split('\n')
        titles = []
        examples = []
        
        for line in lines:
            line = line.strip()
            if line.startswith('Factor') and ':' in line:
                title = line.split(':', 1)[1].strip()
                titles.append(title)
            elif line.startswith('Example:') and ':' in line:
                example = line.split(':', 1)[1].strip()
                examples.append(f"(Ex: {example})")
        
        while len(titles) < 7:
            titles.append(f"Success Factor {len(titles) + 1}")
        while len(examples) < 7:
            examples.append(f"(Ex: Example {len(examples) + 1})")
        
        titles = titles[:7]
        examples = examples[:7]
        
        titles_list = ", ".join(titles)
        
        result = {
            "market_name": context['headline'],
            "titles": titles_list,
            **{f"ksfactor{i+1}_title": titles[i] for i in range(7)},
            **{f"ksfactor{i+1}_example": examples[i] for i in range(7)}
        }
        
        logger.info("Success factors generated successfully")
        return result

    def _generate_market_impacting_factors(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Generate market impacting factors data"""
        logger.info(f"Generating market impacting factors for: {context['headline']}")
        
        market_name = context['headline']
        
        prompt = f"""I want you to act as a Research analyst and provide Top-5 Trending "Market Impacting Factors" for {market_name} with Regional Comparison among (North America, Europe, Asia-Pacific, Latin America, MEA) rated as Low/Medium/High.

Please format your response as follows:
- First provide the headline: "{market_name} Market Impacting Factors"
- Then provide exactly 5 factors, each on a new line starting with "Factor X:"
- Then provide a table with columns: Factor | North America | Europe | Asia Pacific | Latin America | Middle East & Africa
- Each cell in the impact columns should contain only: Low, Medium, or High
- Factor names should be concise (5-7 words)
Example format:
{market_name} Market Impacting Factors

Factor 1: [Factor Name]
Factor 2: [Factor Name] 
Factor 3: [Factor Name]
Factor 4: [Factor Name]
Factor 5: [Factor Name]

| Factor | North America | Europe | Asia Pacific | Latin America | Middle East & Africa |
|--------|---------------|--------|--------------|---------------|---------------------|
| [Factor 1] | High | Medium | Low | Medium | High |
| [Factor 2] | Medium | High | High | Low | Medium |
[etc...]

Please talk about specific market impacting factors relevant to {market_name}."""

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a market research analyst. Provide detailed market impact analysis in the specified format.",
                },
                {"role": "user", "content": prompt},
            ],
        )
        
        ai_response = response.choices[0].message.content
        headline, factors, impact_data = parse_market_factors_response(ai_response)
        
        headline = f"{market_name.upper()} : MARKET IMPACTING FACTORS"
        
        result = {
            "headline": headline,
            "factors": factors,
            "impact_data": impact_data
        }
        
        logger.info("Market impacting factors generated successfully")
        return result

    def _generate_company_info(self, context: Dict[str, Any]) -> Dict[str, str]:
        logger.info(f"Generating company info for: {context['company_name']}")
        prompt = f"""Generate information about {context["company_name"]} in the "{context["headline"]}" domain. 
        Return the information in the following JSON format:
        {{
            "company_name": "{context["company_name"]}",
            "headquarters": "",
            "employee_count": "",
            "revenue": "",
            "top_product": "",
            "description_product": "",
            "estd": "",
            "website": "",
            "geographic_presence": "",
            "ownership": "",
            "short_description_company": ""
        }}
        geographic_presence only choose between from  Global, North America, Europe, Asia Pacific, Latin America, Middle East & Africa
        The short_description_company should be around 100 words. I want you to act as a Research Analyst and give Company Overview of "{context["company_name"]}" in around 10-11 lines (In one paragraph only) which should not talk about Headquarter Country, Establishment/Foundation Year, Number of Employees or Revenue and should not use any marketing/promotional words like, largest, prominent, diversified, recognized, among others (You can talk about its product/service related to {context["headline"]}, market presence, business strategy, recent developments, etc) like this for tone:
        Schlumberger Ltd (SLB) provides technology for reservoir characterization, production, drilling and processing to the oil and gas industry. The company supplies its products and services to the industry, from exploration through production and integrated pipeline solutions for hydrocarbon recovery. SLB's products and services include open-hole and cased-hole wireline logging; drilling services; well completion services, including well testing and artificial lift; well services such as cementing, coiled tubing, stimulations, and sand control; interpretation and consulting services; and integrated project management. The company has an operational presence in North America, Latin America, Europe and Africa, the Middle East and Asia. SLB is headquartered in Houston, Texas, the US..
.       website should be the official website no Https ot http.
        revenue should be in the format " X.XX billion" or " X.XX million" and should be correct 2024 data in USD only correct data must correct strict.
        ownership should be either "Public" or "Private".
        top product should be a product or service relevant to the headline market.
        description_product should be 50 words describing the top product.
        estd is year of establishment should be correct data.
        headquarters should be "Country" format and should be correct data.
        employee_count should be in "X,XXX" or "XX,XXX" format and should be correct data.
        Return ONLY valid JSON, no additional text. no urls/citations for references."""

        response = client.responses.create(
            model="gpt-5",
            tools=[{
                "type": "web_search_preview",
                "search_context_size": "high",
            }],
            input=[
                {"role": "system", "content": "You are a JSON generator. Always return valid JSON and nothing else."},
                {"role": "user", "content": prompt}
            ]
        )          
        content = response.output_text.strip()

        if content.startswith("```json"):
            content = content[7:]
        if content.endswith("```"):
            content = content[:-3]
        content = content.strip()

        result = json.loads(content)
        logger.info("Company info generated successfully")
        return result

    def _generate_research_journals(self, context: Dict[str, Any]) -> List[str]:
        logger.info("Generating research journals")
        market_name = context.get("headline", "Technology Market")

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a JSON generator. Provide the names of research journals related to the specified market from website properly  "
                        "in JSON format. Only include the names as strings, no additional information "
                        "is needed. Search established, reputable journals.\n\n"
                        "Give 5 journal names.\n\n"
                        "**Output format must be a JSON object with a 'journals' key containing an array of strings:**\n"
                        '{"journals": ["Journal Name 1", "Journal Name 2"]}\n'
                        'If there are no journals for the given market, return: {"journals": []}'
                    ),
                },
                {
                    "role": "user",
                    "content": f"Find research journals for: {market_name}",
                },
            ],
            response_format={"type": "json_object"},
        )

        json_response = json.loads(response.choices[0].message.content)
        journals = json_response.get("journals", [])

        default_journals = [
            "Journal of Market Research",
            "International Business Review",
            "Strategic Management Journal",
            "Harvard Business Review",
            "Industrial Marketing Management",
        ]

        if len(journals) < 5:
            journals.extend(default_journals[len(journals) : 5])

        logger.info(f"Research journals generated: {len(journals)} items")
        return journals[:5]

    def _generate_industry_associations(self, context: Dict[str, Any]) -> List[str]:
        logger.info("Generating industry associations")
        market_name = context.get("headline", "Technology Market")

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a JSON generator. Provide the names of industry associations or government organizations "
                        "related to the specified market in JSON format. Only include the names "
                        "as strings, no additional information is needed. Search for highly relevant organizations "
                        "to the market name (exclude private company names). Give 5.\n\n"
                        "**Output format must be a JSON object with an 'associations' key containing an array of strings:**\n"
                        '{"associations": ["Association Name 1", "Association Name 2"]}\n'
                        "If there are no relevant associations or organizations for the given market, "
                        'return: {"associations": []}'
                    ),
                },
                {
                    "role": "user",
                    "content": f"Find industry associations and government organizations for: {market_name}",
                },
            ],
            response_format={"type": "json_object"},
        )

        json_response = json.loads(response.choices[0].message.content)
        associations = json_response.get("associations", [])

        default_associations = [
            "Global Industry Alliance",
            "International Trade Association",
            "National Business Federation",
            "Industry Development Council",
            "Professional Standards Organization",
        ]

        if len(associations) < 5:
            associations.extend(default_associations[len(associations) : 5])

        logger.info(f"Industry associations generated: {len(associations)} items")
        return associations[:5]

    def _generate_financial_overview(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Generate complete financial overview with revenue data, segments, and analysis"""
        logger.info(f"Generating financial overview for: {context['company_name']}")
        company = context["company_name"]

        try:
            # CALL #1: Get revenue data for 2022-2024
            revenue_prompt = f"""
            Search the web and find {company}'s actual annual revenue figures for 2022, 2023, and 2024 in billions USD.
            
            Return ONLY valid JSON in this EXACT format (no text, no explanations):
            {{
                "2022": <number>,
                "2023": <number>,
                "2024": <number>
            }}
            
            Revenue values must be numeric (one digit after decimal, don't round off strictly - example: if it's 49.95 take 49.9 or if its 13.0 take 13.0 ) and represent billions USD.
            Use realistic estimates based on {company}'s actual business scale and recent financial reports.
            Do NOT use null, NA, Unknown, or any non-numeric values.
            When converting it use realtime conversion rate if revenue is in other currency to USD.
            """

            revenue_response = self.openai_client.responses.create(
                model="gpt-5",
                tools=[{
                    "type": "web_search_preview",
                    "search_context_size": "high",
                }],
                input=[
                    {"role": "system", "content": "You are a JSON generator. Always return valid JSON and nothing else."},
                    {"role": "user", "content": revenue_prompt}
                ]
            )          
            revenue_content = revenue_response.output_text.strip()
            if revenue_content.startswith("```json"):
                revenue_content = revenue_content[7:]
            if revenue_content.endswith("```"):
                revenue_content = revenue_content[:-3]
            revenue_content = revenue_content.strip()

            revenue_data = json.loads(revenue_content)
            revenue_data = {
                int(year): float(value) for year, value in revenue_data.items()
            }

            logger.info(
                f"Revenue data fetched: 2022=${revenue_data[2022]}B, 2023=${revenue_data[2023]}B, 2024=${revenue_data[2024]}B"
            )

            segment_prompt = f"""
            Search the web and find ALL main business segments/divisions for {company} along with their 2024 revenue percentage contribution.
            
            Return ONLY valid JSON in this EXACT format:
            {{
                "segments": [
                    {{"name": "Segment1", "percentage": <number>}},
                    {{"name": "Segment2", "percentage": <number>}}
                ]
            }}
            
            Requirements:
            1. Include ALL major revenue-generating segments for {company} based on annual reports
            2. Percentages must add up to approximately 100%
            3. Use actual 2024 business performance data
            4. IF segment not available, provide based on region then only 
            """

            segment_response = self.openai_client.responses.create(
                model="gpt-5",
                tools=[{
                "type": "web_search_preview",
                "search_context_size": "high",
                }],
                input=[
                {"role": "system", "content": "You are a JSON generator. Always return valid JSON and nothing else."},
                {"role": "user", "content": segment_prompt}
                ],
            )

            segment_content = segment_response.output_text.strip()
            if segment_content.startswith("```json"):
                segment_content = segment_content[7:]
            if segment_content.endswith("```"):
                segment_content = segment_content[:-3]
            segment_content = segment_content.strip()

            segment_data = json.loads(segment_content)

            segments_dict = {}
            segment_names = []
            for item in segment_data["segments"]:
                segment_name = item["name"]
                segment_percentage = float(item["percentage"])
                segments_dict[segment_name] = segment_percentage
                segment_names.append(segment_name)

            total = sum(segments_dict.values())
            if not (95 <= total <= 105):
                segments_dict = {
                    seg: (pct / total * 100) for seg, pct in segments_dict.items()
                }

            logger.info(f"Found {len(segment_names)} business segments")

            revenue_analysis = self._generate_revenue_analysis(company, revenue_data)
            segmental_analysis = self._generate_segmental_analysis(
                company, revenue_data, segments_dict
            )

            result = {
                "revenue": revenue_data,
                "segments": segments_dict,
                "segment_names": segment_names,
                "revenue_analysis": revenue_analysis,
                "segmental_analysis": segmental_analysis,
            }

            logger.info("Financial overview generated successfully")
            return result

        except Exception as e:
            logger.error(f"Error generating financial overview: {str(e)}")
            raise

    def _generate_revenue_analysis(self, company: str, revenue_data: Dict) -> str:
        """Generate revenue analysis using gpt-5-mini"""
        prompt = f"""
        Write a professional revenue analysis for {company} based on the following data:
        - 2022: USD {revenue_data[2022]} Billion
        - 2023: USD {revenue_data[2023]} Billion  
        - 2024: USD {revenue_data[2024]} Billion
        
        Requirements:
        1. Write exactly 120-130 words
        2. Focus on year-over-year growth trends and performance
        3. Use professional financial analysis language
        4. One cohesive paragraph format
        5. Include growth percentages and key insights
        """

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a professional financial analyst.",
                },
                {"role": "user", "content": prompt},
            ],
        )

        return response.choices[0].message.content.strip()

    def _generate_segmental_analysis(
        self, company: str, revenue_data: Dict, segments: Dict
    ) -> str:
        """Generate segmental analysis using gpt-5-mini"""
        segment_breakdown = [f"- {seg}: {pct:.1f}%" for seg, pct in segments.items()]
        segments_text = "\n".join(segment_breakdown)

        prompt = f"""
        Write a professional business segmental analysis for {company} based on 2024 data:
        
        Total Revenue 2024: USD {revenue_data[2024]} Billion
        Business Segment Breakdown:
        {segments_text}
        
        Requirements:
        1. Write exactly 120-130 words
        2. Focus on segment performance, opportunities, and market positioning
        3. Use professional business analysis language
        4. One cohesive paragraph format
        5. Highlight key segments and growth opportunities
        """

        response = self.openai_client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a professional business analyst.",
                },
                {"role": "user", "content": prompt},
            ],
        )

        return response.choices[0].message.content.strip()


logger.info("Initializing AI Service globally")
ai_service = AIService(client, api_key_gemini)


class TaxonomyBoxGenerator:
    COLORS = {
        "purple": RGBColor(0x31, 0x09, 0x7E),
        "orange": RGBColor(255, 102, 51),
        "teal": RGBColor(0, 179, 152),
        "blue": RGBColor(0, 162, 232),
        "dark_blue": RGBColor(36, 64, 142),
        "white": RGBColor(255, 255, 255),
        "light_gray": RGBColor(0xF2, 0xF2, 0xF2),
        "text_dark": RGBColor(0, 0, 0),
        "new_blue": RGBColor(0x00, 0x70, 0xC0),
        "light_green": RGBColor(0x92, 0xD0, 0x50),
        "yellow_orange": RGBColor(0xFF, 0xC0, 0x00),
        "dark_red": RGBColor(0xC0, 0x00, 0x00),
        "rose": RGBColor(0xF8, 0x78, 0x84),
        "light_black": RGBColor(0x7F, 0x7F, 0x7F),
        "dark_teal": RGBColor(0x00, 0xA8, 0x8F),
        "turquoise": RGBColor(0x33, 0xC5, 0xF0),
        "new_purple": RGBColor(0x59, 0x46, 0x8F),
    }

    BOX_HEADER_COLORS = [
        COLORS["new_blue"],
        COLORS["light_green"],
        COLORS["yellow_orange"],
        COLORS["dark_red"],
        COLORS["rose"],
        COLORS["light_black"],
        COLORS["dark_teal"],
        COLORS["turquoise"],
        COLORS["new_purple"],
    ]

    def __init__(self, presentation):
        self.prs = presentation
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        self.left_margin, self.top_margin, self.right_margin, self.bottom_margin = (
            Inches(0.5),
            Inches(2),
            Inches(0.5),
            Inches(0.8),
        )
        self.h_spacing, self.v_spacing = Inches(0.2), Inches(0.2)

    def _add_category_box(
        self, slide, category, content, left, top, max_width, max_height, color_index
    ):
        header_color = self.BOX_HEADER_COLORS[color_index % len(self.BOX_HEADER_COLORS)]

        header_height = Inches(0.3)
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max_width, header_height
        )
        header.fill.solid()
        header.fill.fore_color.rgb = header_color
        header.line.color.rgb = header_color
        p = header.text_frame.paragraphs[0]
        p.text = category
        p.font.size, p.font.bold, p.font.color.rgb, p.alignment = (
            Pt(11),
            True,
            self.COLORS["white"],
            PP_ALIGN.CENTER,
        )

        content_box_height = max_height - header_height + Inches(0.2)
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left,
            top + header_height,
            max_width,
            content_box_height,
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = self.COLORS["light_gray"]
        content_box.line.color.rgb = self.COLORS["light_gray"]
        tf = content_box.text_frame
        tf.word_wrap, tf.vertical_anchor = True, MSO_VERTICAL_ANCHOR.TOP
        self._add_list_content(tf, content)

    def _add_list_content(self, text_frame, content):

        text_frame.margin_bottom = Pt(12)
        if text_frame.paragraphs:
            text_frame.paragraphs[0].text = ""

        p = text_frame.paragraphs[0]
        pPr = p._p.get_or_add_pPr()

        lst = pPr.find(qn("a:lstStyle"))
        if lst is None:
            lst = OxmlElement("a:lstStyle")
            pPr.append(lst)

        for i, item in enumerate(content):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = item
            p.alignment = PP_ALIGN.LEFT

            p.line_spacing = 1
            p.space_after = Pt(6)

            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(int(Pt(15).emu)))
            pPr.set("indent", str(int(Pt(-15).emu)))

            marL = OxmlElement("a:marL")
            marL.set("val", str(int(Pt(50).emu)))
            pPr.append(marL)

            indent = OxmlElement("a:indent")
            indent.set("val", str(int(Pt(-19).emu)))
            pPr.append(indent)

            buChar = OxmlElement("a:buChar")
            buChar.set("char", "○")
            pPr.append(buChar)

            buFont = OxmlElement("a:buFont")
            buFont.set("typeface", "Symbol")
            pPr.append(buFont)

            buClr = OxmlElement("a:buClr")
            srgbClr = OxmlElement("a:srgbClr")
            srgbClr.set("val", "000000")
            buClr.append(srgbClr)
            pPr.append(buClr)

            for run in p.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = self.COLORS["text_dark"]

    def add_taxonomy_boxes(self, slide_index, taxonomy_data):
        logger.info(f"Adding taxonomy boxes to slide {slide_index}")
        slide = self.prs.slides[slide_index]
        available_width = self.slide_width - self.left_margin - self.right_margin
        num_categories = len(taxonomy_data)
        boxes_per_row = min(5, num_categories)
        box_width = (
            available_width - (boxes_per_row - 1) * self.h_spacing
        ) / boxes_per_row

        rows, current_row, current_row_width = [], [], 0
        color_index = 0
        for category, hierarchy in taxonomy_data.items():
            item_count = len(hierarchy)
            box_height = max(
                Inches(1), Inches(0.43) + (item_count * Inches(0.22) * 1.2)
            )
            if current_row_width + box_width > available_width and current_row:
                rows.append(current_row)
                current_row, current_row_width = [], 0
            current_row.append(
                {
                    "category": category,
                    "content": hierarchy,
                    "height": box_height,
                    "color_index": color_index,
                }
            )
            current_row_width += box_width + self.h_spacing
            color_index += 1

        if current_row:
            rows.append(current_row)

        current_top = self.top_margin
        for row in rows:
            row_max_height = max(box["height"] for box in row)
            row_width = len(row) * box_width + (len(row) - 1) * self.h_spacing
            left_start = self.left_margin + (available_width - row_width) / 2
            for i, box in enumerate(row):
                left = left_start + i * (box_width + self.h_spacing)
                self._add_category_box(
                    slide,
                    box["category"],
                    box["content"],
                    left,
                    current_top,
                    box_width,
                    box["height"],
                    box["color_index"],
                )
            current_top += row_max_height + self.v_spacing

        logger.info(f"Taxonomy boxes added successfully to slide {slide_index}")


def replace_text_in_presentation(prs, slide_data_dict):
    logger.info("Starting text replacement in presentation")
    for slide_idx, slide in enumerate(prs.slides):
        data = slide_data_dict.get(slide_idx, {})
        if not data:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for key, value in data.items():
                    token = f"{{{{{key}}}}}"
                    if token in p.text:
                        inline_text = p.text
                        p.text = inline_text.replace(token, str(value))
    logger.info("Text replacement in presentation completed")


def replace_text_in_tables(prs, slide_indices, slide_data_dict):
    logger.info(f"Starting text replacement in tables for slides: {slide_indices}")
    for idx in slide_indices:
        if idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        data = slide_data_dict.get(idx, {})
        if not data:
            continue
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            for key, value in data.items():
                                if f"{{{{{key}}}}}" in run.text:
                                    run.text = run.text.replace(
                                        f"{{{{{key}}}}}", str(value)
                                    )
    logger.info("Text replacement in tables completed")


def get_rgb_color_safe(font):
    try:
        return font.color.rgb
    except AttributeError:
        return None


def replace_text_preserving_color(paragraph, placeholder, new_text):
    full_text = "".join(run.text for run in paragraph.runs)

    if placeholder not in full_text:
        return

    for run in paragraph.runs:
        if placeholder in run.text:
            font_color = get_rgb_color_safe(run.font)
            run.text = run.text.replace(placeholder, new_text)
            if font_color:
                run.font.color.rgb = font_color
            break


def replace_text_in_paragraph(paragraph, placeholder, new_text):
    full_text = "".join(run.text for run in paragraph.runs)
    if placeholder in full_text:
        replacement_text = str(new_text) if new_text is not None else ""

        if not replacement_text.strip():
            new_full_text = full_text.replace(placeholder, "").strip()
            import re

            new_full_text = re.sub(r"\n\s*\n", "\n", new_full_text)
        else:
            new_full_text = full_text.replace(placeholder, replacement_text)

        for run in paragraph.runs:
            run.text = ""
        if paragraph.runs:
            paragraph.runs[0].text = new_full_text
        else:
            paragraph.add_run().text = new_full_text


def set_cell_border(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for line_dir in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = OxmlElement(line_dir)
        ln.set("w", "12700")

        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        srgbClr.set("val", "A6A6A6")
        solidFill.append(srgbClr)
        ln.append(solidFill)

        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")

        tcPr.append(ln)


def validate_segment_hierarchy(segment_text):
    logger.info("Validating segment hierarchy")
    lines = segment_text.strip().split("\n")
    errors = []
    last_main_number = 0
    last_sub_numbers = {}

    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue

        match = re.match(r"^(\d+(?:\.\d+)*)\.\s+(.+)$", line)
        if not match:
            errors.append(f"Line {i + 1}: Invalid format")
            continue

        number_parts = [int(n) for n in match.group(1).split(".")]
        depth = len(number_parts)

        if depth == 1:
            if number_parts[0] != last_main_number + 1:
                errors.append(
                    f"Line {i + 1}: Expected main number {last_main_number + 1}, got {number_parts[0]}"
                )
            last_main_number = number_parts[0]
            last_sub_numbers = {}
        elif depth == 2:
            main_num = number_parts[0]
            sub_num = number_parts[1]

            if main_num != last_main_number:
                errors.append(
                    f"Line {i + 1}: Sub-item doesn't match current main number"
                )

            expected_sub = last_sub_numbers.get(main_num, 0) + 1
            if sub_num != expected_sub:
                errors.append(
                    f"Line {i + 1}: Expected sub-number {main_num}.{expected_sub}"
                )
            last_sub_numbers[main_num] = sub_num
        elif depth == 3:
            main_num = number_parts[0]
            sub_num = number_parts[1]
            sub_sub_num = number_parts[2]

            if main_num != last_main_number:
                errors.append(
                    f"Line {i + 1}: Sub-sub-item doesn't match current main number"
                )

            key = f"{main_num}.{sub_num}"
            expected_sub_sub = last_sub_numbers.get(key, 0) + 1
            if sub_sub_num != expected_sub_sub:
                errors.append(
                    f"Line {i + 1}: Expected sub-sub-number {main_num}.{sub_num}.{expected_sub_sub}"
                )
            last_sub_numbers[key] = sub_sub_num

    if errors:
        logger.warning(f"Segment hierarchy validation found {len(errors)} errors")
    else:
        logger.info("Segment hierarchy validation passed")
    return errors


def generate_actual_data():
    data = [
        [2019, 7.0, 5.0, 4.0, 3.5, 3.1, 2.5, 2.0],
        [2020, 7.5, 5.4, 4.3, 3.7, 3.3, 2.7, 2.2],
        [2021, 8.1, 5.7, 4.6, 4.0, 3.6, 2.9, 2.3],
        [2022, 8.7, 6.1, 5.0, 4.3, 3.8, 3.1, 2.5],
        [2023, 9.3, 6.6, 5.3, 4.7, 4.1, 3.3, 2.7],
        [2024, 9.9, 7.1, 5.8, 5.0, 4.4, 3.6, 2.9],
        [2025, 10.7, 7.7, 6.2, 5.4, 4.8, 3.8, 3.1],
        [2026, 11.6, 8.2, 6.7, 5.8, 5.1, 4.1, 3.3],
        [2027, 12.4, 8.9, 7.2, 6.2, 5.5, 4.5, 3.6],
        [2028, 13.4, 9.6, 7.7, 6.7, 5.9, 4.8, 3.8],
        [2029, 14.4, 10.3, 8.4, 7.2, 6.4, 5.1, 4.2],
        [2030, 15.4, 11.1, 9.0, 7.7, 6.8, 5.6, 4.5],
        [2031, 16.6, 11.9, 9.7, 8.4, 7.3, 5.9, 4.8],
        [2032, 17.8, 12.8, 10.4, 9.0, 7.9, 6.4, 5.2],
    ]
    return data


def parse_segment_input(segment_input: str) -> Dict[str, Dict]:
    logger.info("Parsing segment input")
    lines = segment_input.strip().split("\n")
    nested_dict = {}
    level_stack = []
    for line in lines:
        if not line.strip():
            continue
        
        match = re.match(r'^(\d+(?:\.\d+)*)\.\s+(.+)$', line.strip())
        if not match:
            continue
            
        key = match.group(1)
        value = match.group(2).strip()

        parts = key.split(".")
        depth = len(parts)
        label = value
        level_stack = level_stack[: depth - 1]
        current = nested_dict
        for k in level_stack:
            current = current[k]
        current[label] = {}
        level_stack.append(label)
    logger.info(
        f"Segment input parsed successfully: {len(nested_dict)} main categories"
    )
    return nested_dict


def generate_toc_data(
    nested_dict: Dict,
    headline: str,
    forecast_period: str,
    user_segment: str,
    kmi_items: List[str] = None,
) -> Dict[str, int]:
    logger.info("Generating Table of Contents data")
    toc_start_levels = {
        "1. Introduction": 0,
        "1.1. Objectives of the Study": 1,
        "1.2. Market Definition & Scope": 1,
        "2. Research Methodology": 0,
        "2.1. Research Process": 1,
        "2.2. Secondary & Primary Data Methods": 1,
        "2.3. Market Size Estimation Methods": 1,
        "2.4. Market Assumptions & Limitations": 1,
        "3. Executive Summary": 0,
        "3.1. Global Market Outlook": 1,
        "3.2. Key Market Highlights": 1,
        "3.3. Segmental Overview": 1,
        "4. Market Dynamics & Outlook": 0,
        "4.1. Macro-Economic Indicators​": 1,
        "4.2. Drivers & Opportunities": 1,
        "4.3. Restraints & Challenges": 1,
        "4.4. Supply Side Trends": 1,
        "4.5. Demand Side Trends": 1,
        "4.6. Porter's Analysis & Impact": 1,
        "4.6.1. Competitive Rivalry": 2,
        "4.6.2. Threat of substitutes": 2,
        "4.6.3. Bargaining power of buyers": 2,
        "4.6.4. Threat of new entrants": 2,
        "4.6.5. Bargaining power of suppliers": 2,
    }

    kmi_section = {"5. Key Market Insights": 0}

    default_kmi_items = [
        "Key Success Factors",
        "Market Impacting Factors",
        "Top Investment Pockets",
        "Market Attractiveness Index, 2024",
        "Market Ecosystem",
        "PESTEL Analysis",
        "Pricing Analysis",
        "Regulatory Landscape",
    ]

    all_kmi_items = default_kmi_items.copy()
    if kmi_items:
        all_kmi_items.extend(kmi_items)

    for i, kmi_item in enumerate(all_kmi_items, start=1):
        kmi_section[f"5.{i}. {kmi_item}"] = 1

    def add_nested_items(items_dict, prefix, base_level):
        """Recursively add nested items to TOC with proper levels"""
        item_count = 1
        for item_name, sub_items in items_dict.items():
            current_key = f"{prefix}.{item_count}. {item_name}"
            toc_mid[current_key] = base_level

            if sub_items and isinstance(sub_items, dict):
                add_nested_items(sub_items, f"{prefix}.{item_count}", base_level + 1)

            item_count += 1

    toc_mid = {}
    main_index = 6
    for type_index, (type_name, points) in enumerate(
        nested_dict.items(), start=main_index
    ):
        toc_mid[f"{type_index}. {headline} Size by {type_name} (2019-2032)"] = 0
        add_nested_items(points, str(type_index), 1)

    x = len(list(nested_dict.keys())) + 6
    toc_end_levels = {
        f"{x}. {headline} Size by Region (2019-2032)": 0,
        f"{x}.1. North America ({user_segment})": 1,
        f"{x}.1.1. US": 2,
        f"{x}.1.2. Canada": 2,
        f"{x}.2. Europe ({user_segment})": 1,
        f"{x}.2.1. UK": 2,
        f"{x}.2.2. Germany": 2,
        f"{x}.2.3. Spain": 2,
        f"{x}.2.4. France": 2,
        f"{x}.2.5. Italy": 2,
        f"{x}.2.6. Rest of Europe": 2,
        f"{x}.3. Asia-Pacific ({user_segment})": 1,
        f"{x}.3.1. China": 2,
        f"{x}.3.2. India": 2,
        f"{x}.3.3. Japan": 2,
        f"{x}.3.4. South Korea": 2,
        f"{x}.3.5. Rest of Asia Pacific": 2,
        f"{x}.4. Latin America ({user_segment})": 1,
        f"{x}.4.1. Brazil": 2,
        f"{x}.4.2. Mexico": 2,
        f"{x}.4.3. Rest of Latin America": 2,
        f"{x}.5. Middle East & Africa ({user_segment})": 1,
        f"{x}.5.1. GCC Countries": 2,
        f"{x}.5.2. South Africa": 2,
        f"{x}.5.3. Rest of Middle East & Africa": 2,
        f"{x+1}. Competitive Landscape": 0,
        f"{x+1}.1. Competitive Dashboard": 1,
        f"{x+1}.2. Market Positioning of Key Players, 2024": 1,
        f"{x+1}.3. Strategies Adopted by Key Market Players": 1,
        f"{x+1}.4. Recent Developments in the Market": 1,
        f"{x+1}.5. Company Market Share Analysis, 2024": 1,
        f"{x+2}. Key Company Profiles": 0,
    }

    logger.info(
        f"TOC data generated with {len(toc_start_levels) + len(kmi_section) + len(toc_mid) + len(toc_end_levels)} items"
    )
    return {**toc_start_levels, **kmi_section, **toc_mid, **toc_end_levels}


def add_toc_to_slides(
    prs: Presentation, toc_data_levels: Dict[str, int], toc_slide_indices: List[int]
):
    logger.info(f"Adding TOC to slides: {toc_slide_indices}")
    for i in toc_slide_indices:
        slide = prs.slides[i]
        table_shape = slide.shapes.add_table(
            17, 2, Inches(2.8), Inches(0.5), Inches(10), Inches(6)
        )
        table = table_shape.table
        for row in table.rows:
            for cell in row.cells:
                cell.text = ""
                cell.fill.background()
                tcPr = cell._tc.get_or_add_tcPr()
                for border_tag in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
                    tcPr.append(
                        parse_xml(
                            f'<{border_tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:noFill/></{border_tag}>'
                        )
                    )

    content_items = list(toc_data_levels.keys())
    content_index = 0
    for i in toc_slide_indices:
        table = prs.slides[i].shapes[-1].table
        for col in range(2):
            for row in range(17):
                if content_index >= len(content_items):
                    break
                cell, key = table.cell(row, col), content_items[content_index]
                level = toc_data_levels[key]
                para = cell.text_frame.paragraphs[0]
                para.text = "          " * level + key
                font = para.font
                font.color.rgb, font.size, font.name = (
                    RGBColor(0, 0, 0),
                    Pt(11),
                    "Poppins",
                )
                if key.startswith("The following companies") or key.startswith(
                    "Note :"
                ):
                    font.size = Pt(9)
                    font.color.rgb, font.bold = RGBColor(112, 48, 160), True
                else:
                    font.size = Pt(11)
                if (
                    level == 0
                    and not key.startswith("The following companies")
                    and not key.startswith("Note :")
                ):
                    font.color.rgb, font.bold = RGBColor(112, 48, 160), True
                else:
                    font.color.rgb = RGBColor(0, 0, 0)
                    font.bold = False
                content_index += 1
    logger.info("TOC added to slides successfully")


def create_chart_on_slide(
    slide: Any,
    data: List[List],
    chart_columns: List[str],
    left: float,
    top: float,
    width: float,
    height: float,
):
    logger.info(
        f"Creating chart with {len(chart_columns)} series and {len(data)} data points"
    )
    chart_data = CategoryChartData()
    chart_data.categories = [str(row[0]) for row in data]

    num_series = min(len(chart_columns), 7)
    for i in range(num_series):
        chart_data.add_series(chart_columns[i], [row[i + 1] for row in data])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, chart_data
    ).chart

    chart.plots[0].gap_width = 150
    chart.chart_style = 2
    chart.has_title = False

    chart.has_legend = True
    chart.legend.font.size = Pt(8)
    chart.legend.font.name = "Poppins"
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    value_axis = chart.value_axis
    value_axis.visible = False
    value_axis.has_major_gridlines = False
    value_axis.has_minor_gridlines = False
    value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE

    category_axis = chart.category_axis
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.minor_tick_mark = XL_TICK_MARK.NONE
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(10)
    cat_axis.tick_labels.font.name = "Poppins"
    cat_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    logger.info("Chart created successfully")
def create_financial_charts_on_slide(slide, financial_data, company):
    """Create bar chart and pie chart for financial overview on slide 36"""
    logger.info(f"Creating financial charts for {company}")

    chart_data = CategoryChartData()
    chart_data.categories = ["2022", "2023", "2024"]
    chart_data.add_series(
        "Revenue (USD Bn)",
        [
            financial_data["revenue"][2022],
            financial_data["revenue"][2023],
            financial_data["revenue"][2024],
        ],
    )

    bar_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.5),
        Inches(1.5),
        Inches(4.2),
        Inches(2.5),
        chart_data,
    ).chart

    bar_chart.has_title = False
    bar_chart.has_legend = False
    bar_chart.plots[0].gap_width = 50
    bar_chart.plots[0].has_data_labels = True

    data_labels = bar_chart.plots[0].data_labels
    data_labels.font.size = Pt(10)
    data_labels.font.name = "Poppins"
    data_labels.font.bold = False

    value_axis = bar_chart.value_axis
    value_axis.visible = False
    value_axis.has_major_gridlines = False
    value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE

    category_axis = bar_chart.category_axis
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.tick_labels.font.size = Pt(10)
    category_axis.tick_labels.font.name = "Poppins"

    series = bar_chart.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(49, 9, 126)

    logger.info("Revenue bar chart created")

    pie_data = CategoryChartData()
    pie_data.categories = financial_data["segment_names"]
    pie_data.add_series("Share (%)", list(financial_data["segments"].values()))

    pie_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(7.0), Inches(1.2), Inches(4.5), Inches(3), pie_data
    ).chart

    num_segments = len(financial_data["segment_names"])
    legend_font_size = max(7, min(10, 10 - (num_segments - 4)))

    pie_chart.has_title = False
    pie_chart.has_legend = True
    pie_chart.legend.position = XL_LEGEND_POSITION.RIGHT
    pie_chart.legend.font.size = Pt(legend_font_size)
    pie_chart.legend.font.name = "Poppins"
    pie_chart.legend.include_in_layout = False
    pie_chart.legend.overlay = False

    pie_chart.plots[0].has_data_labels = True
    data_labels = pie_chart.plots[0].data_labels
    data_labels.font.size = Pt(max(7, 10 - (num_segments - 4)))
    data_labels.font.name = "Poppins"
    data_labels.show_value = False
    data_labels.show_percentage = True
    data_labels.show_category_name = False

    series = pie_chart.series[0]
    points = series.points

    color_palette = [
        RGBColor(49, 9, 126),
        RGBColor(51, 197, 240),
        RGBColor(255, 165, 0),
        RGBColor(34, 139, 34),
        RGBColor(220, 20, 60),
        RGBColor(255, 215, 0),
        RGBColor(138, 43, 226),
        RGBColor(0, 191, 255),
        RGBColor(255, 105, 180),
        RGBColor(46, 139, 87),
    ]

    for i, point in enumerate(points):
        color_index = i % len(color_palette)
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color_palette[color_index]

    logger.info(f"Pie chart created with {num_segments} segments")


def clean_filename(filename):
    logger.info(f"Cleaning filename: {filename}")
    invalid_chars = '<>:"/\\|?*'
    for char in invalid_chars:
        filename = filename.replace(char, "_")

    filename = " ".join(filename.split())
    filename = filename[:100]

    logger.info(f"Cleaned filename: {filename}")
    return filename


@app.route("/")
@app.route("/gii/")
def index():
    logger.info("Index page accessed")
    return render_template("index.html")


@app.route("/gii/generate-ppt", methods=["POST"])
def generate_ppt():
    start_time = time.time()
    request_id = f"req_{int(time.time())}"
    logger.info(f"[{request_id}] PPT generation request started")

    try:
        form_data = request.form
        required_fields = [
            "headline",
            "headline_2",
            "historical_year",
            "base_year",
            "forecast_year",
            "forecast_period",
            "cur",
            "value_in",
            "rev_current",
            "rev_future",
            "segment_input",
            "companies",
            "cagr",
            "sqcode",
        ]

        missing_fields = []
        for field in required_fields:
            if not form_data.get(field, "").strip():
                missing_fields.append(field)

        if missing_fields:
            logger.warning(f"[{request_id}] Missing required fields: {missing_fields}")
            return (
                jsonify({"error": "Missing required fields", "fields": missing_fields}),
                400,
            )

        segment_errors = validate_segment_hierarchy(form_data["segment_input"])
        if segment_errors:
            logger.warning(
                f"[{request_id}] Segment hierarchy validation failed: {len(segment_errors)} errors"
            )
            return (
                jsonify(
                    {"error": "Invalid segment hierarchy", "details": segment_errors}
                ),
                400,
            )

        logger.info(f"[{request_id}] Form validation passed")

        headline = form_data["headline"]
        headline_2 = headline.upper()
        headline_3 = headline_2.replace("GLOBAL", "").strip()
        historical_year = "2019-2023"
        base_year = "2024"
        forecast_year = "2032"
        forecast_period = "2025-2032"
        cur = "USD"
        value_in = form_data["value_in"]
        currency = f"{cur} {value_in}"
        rev_current = form_data["rev_current"]
        rev_future = form_data["rev_future"]
        cagr = form_data.get("cagr")
        segment_input = form_data["segment_input"]
        sqcode = form_data.get("sqcode", "").strip()
        kmi_items = []
        kmi_input = form_data.get("kmi_items", "").strip()

        if kmi_input:
            kmi_items = [item.strip() for item in kmi_input.split("\n") if item.strip()]
            logger.info(f"[{request_id}] Custom KMI items provided: {len(kmi_items)}")

        def format_as_bullets(items_list):
            if not items_list:
                return ""
            return "\n".join([f"{item}" for item in items_list])

        default_kmiitems = [
            "Key Success Factors",
            "Market Impacting Factors",
            "Top Investment Pockets",
            "Market Attractiveness Index, 2024",
            "Market Ecosystem",
            "PESTEL Analysis",
            "Pricing Analysis",
            "Regulatory Landscape",
        ]
        default_kmi_bullets = format_as_bullets(default_kmiitems)
        user_kmi_bullets = format_as_bullets(kmi_items) if kmi_items else ""

        companies_input = form_data["companies"].strip()
        company_list = [
            company.strip()
            for company in companies_input.split("\n")
            if company.strip()
        ]

        if not company_list:
            logger.warning(f"[{request_id}] No companies provided")
            return (
                jsonify(
                    {
                        "error": "At least one company must be provided",
                        "message": "Please provide company names, one per line",
                    }
                ),
                400,
            )

        logger.info(f"[{request_id}] Processing {len(company_list)} companies")

        nested_dict = parse_segment_input(segment_input)
        main_topic = list(nested_dict.keys())
        s_segment = "By " + "\nBy ".join(main_topic)
        user_segment = "By " + ", By ".join(main_topic)

        output_lines = []
        for main_type, points in nested_dict.items():
            line_parts = []
            for point, subpoints in points.items():
                if subpoints:
                    subpoint_str = ", ".join(subpoints.keys())
                    line_parts.append(f"{point} ({subpoint_str})")
                else:
                    line_parts.append(point)
            output_lines.append(f"By {main_type}: {', '.join(line_parts)}")
        output_lines.append(
            "By Region: North America, Europe, Asia-Pacific, Latin America, Middle East & Africa"
        )
        context = "\n".join(output_lines)
        logger.info(f"[{request_id}] Context generated successfully")

        toc_data_levels = generate_toc_data(
            nested_dict, headline, forecast_period, user_segment, kmi_items
        )

        ai_context = {
            "headline": headline,
            "value_in": value_in,
            "cur": cur,
            "historical_year": historical_year,
            "forecast_year": forecast_year,
            "base_year": base_year,
            "rev_current": rev_current,
            "rev_future": rev_future,
            "main_topic": main_topic[0] if main_topic else "Type 1",
            "currency": currency.upper(),
            "cagr": cagr,
            "company_name": company_list[0],
        }

        logger.info(f"[{request_id}] Starting AI content generation")
        ai_start_time = time.time()
        ai_results = ai_service.generate_content_parallel(ai_context, nested_dict, headline)        
        ai_elapsed = time.time() - ai_start_time
        timing_logger.info(
            f"[{request_id}] AI content generation completed in {ai_elapsed:.2f}s"
        )

        mpara_11 = ai_results["executive_summary"]
        para_11 = ai_results["market_enablers"]
        para_14_dict = ai_results["industry_expansion"]
        industry_title = para_14_dict["title"]
        para_14_dict_1 = ai_results["industry_expansion_1"]
        para_15_dict = ai_results["investment_challenges"]
        research_journals = ai_results["research_journals"]
        industry_associations = ai_results["industry_associations"]
        company_info = ai_results["company_info"]
        financial_overview = ai_results["financial_overview"]
        success_factors_data = ai_results["success_factors"]
        market_factors_data = ai_results["market_impacting_factors"]
        industry_title_1 = para_15_dict["title"]
        industry_title_2 = para_14_dict_1["title"]
        para_14_1 = "\n".join(para_14_dict_1["paragraphs"])
        para_15 = "\n".join(para_15_dict["paragraphs"])
        para_14 = "\n".join(para_14_dict["paragraphs"])
        definitions = ai_results["definitions"]

        logger.info(f"[{request_id}] AI content extracted successfully")

        x = len(main_topic) + 6

        first_company = company_list[0]
        toc_data_levels[f"{x+2}.1. {first_company}"] = 1
        toc_data_levels[f"{x+2}.1.1. Company Overview"] = 2
        toc_data_levels[f"{x+2}.1.2. Product Portfolio Overview"] = 2
        toc_data_levels[f"{x+2}.1.3. Financial Overview"] = 2
        toc_data_levels[f"{x+2}.1.4. Key Developments"] = 2

        toc_data_levels[
            "The following companies are listed for indicative purposes only. Similar information will be provided for each, with detailed financial data available exclusively for publicly listed companies."
        ] = 0
        for i, name in enumerate(company_list[1:], start=2):
            toc_data_levels[f"{x+2}.{i}. {name}"] = 1
        toc_data_levels[
            "Note : The list of companies mentioned are for indication purpose and are subject to change over the due course of the research"
        ] = 0

        toc_data_levels[f"{x+3}. Conclusion & Recommendation"] = 0
        table_taxonomy = {
            f"BY {key.upper()}": list(value.keys())
            for key, value in nested_dict.items()
        }

        logger.info(f"[{request_id}] Starting presentation modification")
        ppt_start_time = time.time()

        template_path = "testpptgii.pptx"
        if not os.path.exists(template_path):
            logger.error(f"[{request_id}] Template file not found: {template_path}")
            return (
                jsonify(
                    {
                        "error": "Template file not found",
                        "message": "Please ensure testpptgii.pptx is in the project directory",
                    }
                ),
                500,
            )

        logger.info(f"[{request_id}] Loading base presentation")
        prs = Presentation(template_path)

        # ============ INSERT DEFINITION TABLES AT SLIDE 10 ============
        DEFINITION_START_INDEX = 10
        logger.info(
            f"[{request_id}] Inserting definition tables at slide {DEFINITION_START_INDEX}"
        )
        num_definition_slides = insert_definition_tables(
            prs, nested_dict, headline, DEFINITION_START_INDEX
        )
        logger.info(
            f"[{request_id}] Definition slides created: {num_definition_slides}"
        )

        # Calculate slide shift
        SLIDE_SHIFT = num_definition_slides
        logger.info(f"[{request_id}] Slide shift amount: {SLIDE_SHIFT}")

        # ============ APPLY SUCCESS FACTORS (25 + SLIDE_SHIFT) ============
        success_factors_slide_index = 25 + SLIDE_SHIFT
        logger.info(f"[{request_id}] Applying success factors to slide {success_factors_slide_index}")
        
        if success_factors_slide_index < len(prs.slides):
            success_factors_slide = prs.slides[success_factors_slide_index]
            apply_success_factors_to_slide(success_factors_slide, success_factors_data)
            logger.info(f"[{request_id}] Success factors applied successfully")
        else:
            logger.warning(f"[{request_id}] Success factors slide index {success_factors_slide_index} out of range")

        # ============ ADJUST ALL SLIDE DATA INDICES ============
        slide_data = {
            0: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
                "context": context,
            },
            1: {
                "heading_2": f"{headline_2} ({currency.upper()})",
                "hyear": f"Historical Year - {historical_year}",
                "fyear": f"Forecast Year - {forecast_period}",
                "byear": f"Base Year - {base_year}",
            },
            2: {
                "heading_2": f"{headline_2} ({currency.upper()})",
                "hyear": f"Historical Year - {historical_year}",
                "fyear": f"Forecast Year - {forecast_period}",
                "byear": f"Base Year - {base_year}",
            },
            3: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
            },
            9: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
            },
            10 + SLIDE_SHIFT: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
            },
            12 + SLIDE_SHIFT: {
                "org_1": (
                    industry_associations[0]
                    if len(industry_associations) > 0
                    else "Global Industry Alliance"
                ),
                "org_2": (
                    industry_associations[1]
                    if len(industry_associations) > 1
                    else "International Trade Association"
                ),
                "org_3": (
                    industry_associations[2]
                    if len(industry_associations) > 2
                    else "National Business Federation"
                ),
                "org_4": (
                    industry_associations[3]
                    if len(industry_associations) > 3
                    else "Industry Development Council"
                ),
                "org_5": (
                    industry_associations[4]
                    if len(industry_associations) > 4
                    else "Professional Standards Organization"
                ),
                "paper_1": (
                    research_journals[0]
                    if len(research_journals) > 0
                    else "Journal of Market Research"
                ),
                "paper_2": (
                    research_journals[1]
                    if len(research_journals) > 1
                    else "International Business Review"
                ),
                "paper_3": (
                    research_journals[2]
                    if len(research_journals) > 2
                    else "Strategic Management Journal"
                ),
                "paper_4": (
                    research_journals[3]
                    if len(research_journals) > 3
                    else "Harvard Business Review"
                ),
                "paper_5": (
                    research_journals[4]
                    if len(research_journals) > 4
                    else "Industrial Marketing Management"
                ),
            },
            14 + SLIDE_SHIFT: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
            },
            15 + SLIDE_SHIFT: {
                "heading_2": f"{headline_2} SIZE, ({currency.upper()})",
                "mpara": mpara_11,
                "para": para_11,
                "amount_1": rev_current,
                "amount_2": rev_future,
            },
            16 + SLIDE_SHIFT: {
                "heading": headline_2,
                "amount_1": f"{cur} {rev_current} {value_in.upper()} ",
                "amount_2": f"{rev_future} {value_in.upper()} {cur}",
            },
            17 + SLIDE_SHIFT: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
            },
            19 + SLIDE_SHIFT: {"industry_title": industry_title, "para": para_14},
            20 + SLIDE_SHIFT: {"industry_title": industry_title_2, "para": para_14_1},
            21 + SLIDE_SHIFT: {"industry_title": industry_title_1, "para": para_15},
            23 + SLIDE_SHIFT: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
                "default_kmi": default_kmi_bullets,
                "user_kmi": user_kmi_bullets,
            },
            # Slide 24 + SLIDE_SHIFT is Market Impacting Factors (processed separately)
            # Slide 25 + SLIDE_SHIFT is Success Factors (processed above)
            27 + SLIDE_SHIFT: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
                "types": s_segment,
            },
            28 + SLIDE_SHIFT: {
                "heading": headline_2,
                "type_1": main_topic[0].upper() if main_topic else "Type 1",
                "timeline": "2019-2032",
                "cur": f"{cur.upper()} {value_in.upper()}",
            },
            29 + SLIDE_SHIFT: {
                "heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
            },
            30 + SLIDE_SHIFT: {
                "heading": headline_2,
                "timeline": "2019-2032",
                "cur": f"{cur.upper()} {value_in.upper()}",
            },
            31 + SLIDE_SHIFT: {
                "2_heading": headline_3.upper(),
                "type_1": main_topic[0].upper() if main_topic else "Type 1",
                "timeline": "2019-2032",
                "cur": f"{cur.upper()} {value_in.upper()}",
            },
            32 + SLIDE_SHIFT: {
                "2_heading": headline_3.upper(),
                "type_1": main_topic[0].upper() if main_topic else "Type 1",
                "timeline": "2019-2032",
                "cur": f"{cur.upper()} {value_in.upper()}",
            },
            33 + SLIDE_SHIFT: {
                "2_heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
            },
            34 + SLIDE_SHIFT: {"heading": headline},
            35 + SLIDE_SHIFT: {
                "company": company_info["company_name"].upper(),
                "e": company_info["employee_count"],
                "h": company_info["headquarters"],
                "geo": company_info["geographic_presence"],
                "es": company_info["estd"],
                "rev": company_info["revenue"],
            },
            36 + SLIDE_SHIFT: {
                "2_heading": headline_2,
                "timeline": f"HISTORIC YEAR {historical_year} \nFORECAST TO {forecast_year}",
            },
            37 + SLIDE_SHIFT: {
                "company": company_info["company_name"].upper(),
                "e": company_info["employee_count"],
                "ownership": company_info["ownership"],
                "h": company_info["headquarters"],
                "website": company_info["website"],
                "es": company_info["estd"],
                "product": company_info["top_product"],
                "para": company_info["short_description_company"],
                "rev": company_info["revenue"],
                "geo": company_info["geographic_presence"],
                "description": company_info["description_product"],
            },
            38 + SLIDE_SHIFT: {
                "company": company_info["company_name"].upper(),
                "revenue_analysis": financial_overview["revenue_analysis"],
                "segmental_analysis": financial_overview["segmental_analysis"],
            },
            39 + SLIDE_SHIFT: {"company": company_info["company_name"].upper()},
        }

        # Process regular slide data
        for slide in prs.slides:
            data = slide_data.get(prs.slides.index(slide), {})
            if not data:
                continue
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for key, value in data.items():
                            token = f"{{{{{key}}}}}"
                            replace_text_in_paragraph(paragraph, token, value)

        # Second pass for color preservation
        for slide in prs.slides:
            data = slide_data.get(prs.slides.index(slide), {})
            if not data:
                continue
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for key, value in data.items():
                            token = f"{{{{{key}}}}}"
                            replace_text_preserving_color(paragraph, token, value)

        # ============ APPLY MARKET IMPACTING FACTORS (24 + SLIDE_SHIFT) ============
        logger.info(f"[{request_id}] Applying market impacting factors to slide")
        market_factors_slide_index = 24 + SLIDE_SHIFT
        apply_market_factors_to_slide(prs, market_factors_slide_index, market_factors_data)

        logger.info(f"[{request_id}] Adding taxonomy boxes")
        generator = TaxonomyBoxGenerator(prs)
        generator.add_taxonomy_boxes(1, table_taxonomy)

        logger.info(f"[{request_id}] Performing text replacements in tables")
        table_slide_indices = [
            12 + SLIDE_SHIFT,
            15 + SLIDE_SHIFT,
            18 + SLIDE_SHIFT,
            19 + SLIDE_SHIFT,
            20 + SLIDE_SHIFT,
            21 + SLIDE_SHIFT,
            23 + SLIDE_SHIFT,
            24 + SLIDE_SHIFT,
            25 + SLIDE_SHIFT,  # Success Factors slide
            27 + SLIDE_SHIFT,
            28 + SLIDE_SHIFT,
            29 + SLIDE_SHIFT,
            30 + SLIDE_SHIFT,
            31 + SLIDE_SHIFT,
            32 + SLIDE_SHIFT,
            33 + SLIDE_SHIFT,
            34 + SLIDE_SHIFT,
            35 + SLIDE_SHIFT,
            36 + SLIDE_SHIFT,
            37 + SLIDE_SHIFT,
            38 + SLIDE_SHIFT,
            39 + SLIDE_SHIFT,
        ]
        replace_text_in_tables(prs, table_slide_indices, slide_data)

        logger.info(f"[{request_id}] Creating Table of Contents")
        toc_slide_indices = [4, 5, 6, 7, 8]
        add_toc_to_slides(prs, toc_data_levels, toc_slide_indices)

        logger.info(f"[{request_id}] Adding tables and charts")
        target_slide_indices = [28 + SLIDE_SHIFT, 31 + SLIDE_SHIFT, 32 + SLIDE_SHIFT]
        graph_table = list(nested_dict[main_topic[0]].keys()) if main_topic else []
        total_rows = len(graph_table)

        row_labels = graph_table.copy()
        row_labels.append("Total")

        years = [str(y) for y in range(2019, 2033)]
        columns = [""] + years + ["CAGR (2025–2032)"]
        num_rows = len(row_labels) + 1
        num_cols = len(columns)

        header_rgb = RGBColor(49, 6, 126)
        border_rgb = RGBColor(166, 166, 166)
        alt_row_colors = [RGBColor(231, 231, 231), RGBColor(255, 255, 255)]

        font_mapping = {
            "header": "Poppins Bold",
            "first_col": "Poppins Bold",
            "values": "Poppins Medium",
        }

        for slide_index in target_slide_indices:
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]

                left = Inches(0.45)
                top = Inches(4.05)
                width = Inches(8.7)
                height = Inches(0.72 + num_rows * 0.3)
                table = slide.shapes.add_table(
                    num_rows, num_cols, left, top, width, height
                ).table

                for col_index, header in enumerate(columns):
                    cell = table.cell(0, col_index)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_rgb

                    cell.text_frame.clear()
                    para = cell.text_frame.paragraphs[0]
                    para.text = header.replace("\n", " ").strip()
                    para.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                    if para.runs:
                        run = para.runs[0]
                    else:
                        run = para.add_run()

                    if col_index != num_cols - 1:
                        run.font.size = Pt(5.7)
                        cell.text_frame.word_wrap = False
                    else:
                        run.font.size = Pt(8)

                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.name = font_mapping["header"]

                for row_index, label in enumerate(row_labels, start=1):
                    row_color = alt_row_colors[(row_index - 1) % 2]

                    for col_index in range(num_cols):
                        cell = table.cell(row_index, col_index)

                        if col_index == 0:
                            cell.text = label
                        elif col_index == num_cols - 1:
                            cell.text = "XX%"
                        else:
                            cell.text = "XX"

                        para = cell.text_frame.paragraphs[0]
                        para.alignment = PP_ALIGN.CENTER
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                        if col_index == 0:
                            para.font.size = Pt(8)
                            para.font.name = font_mapping["first_col"]
                            para.font.bold = True
                        else:
                            para.font.size = Pt(9)
                            para.font.name = font_mapping["values"]

                            if label == "Total" and col_index == num_cols - 1:
                                para.font.bold = True
                            if row_index == num_rows - 1:
                                para.font.bold = True

                        cell.fill.solid()
                        cell.fill.fore_color.rgb = row_color
                        set_cell_border(cell)

                for col_index in range(num_cols):
                    if col_index == 0:
                        table.columns[col_index].width = Inches(1)
                    elif col_index == num_cols - 1:
                        table.columns[col_index].width = Inches(0.8)
                    else:
                        table.columns[col_index].width = Inches(0.4)

        if main_topic:
            chart_columns = graph_table

            for idx in target_slide_indices:
                if idx < len(prs.slides):
                    slide = prs.slides[idx]
                    data = generate_actual_data()
                    create_chart_on_slide(
                        slide,
                        data,
                        chart_columns,
                        Inches(0.4),
                        Inches(1.1),
                        Inches(12.5),
                        Inches(2.8),
                    )

        logger.info(
            f"[{request_id}] Adding financial charts to slide {38 + SLIDE_SHIFT}"
        )
        slide_38 = prs.slides[38 + SLIDE_SHIFT]
        create_financial_charts_on_slide(
            slide_38, financial_overview, company_info["company_name"]
        )

        data_38 = slide_data.get(38 + SLIDE_SHIFT, {})

        for shape in slide_38.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = "".join(run.text for run in paragraph.runs)
                    for key, value in data_38.items():
                        token = f"{{{{{key}}}}}"
                        if token in full_text:
                            new_text = full_text.replace(token, str(value))
                            paragraph.text = new_text
                            for run in paragraph.runs:
                                run.font.size = Pt(10)
                                run.font.name = "Poppins"
                                run.font.bold = False
                                run.font.color.rgb = RGBColor(0, 0, 0)

            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            full_text = "".join(run.text for run in paragraph.runs)
                            for key, value in data_38.items():
                                token = f"{{{{{key}}}}}"
                                if token in full_text:
                                    new_text = full_text.replace(token, str(value))
                                    paragraph.text = new_text
                                    for run in paragraph.runs:
                                        run.font.size = Pt(10)
                                        run.font.name = "Poppins"
                                        run.font.bold = False
                                        run.font.color.rgb = RGBColor(0, 0, 0)

        clean_market_name = clean_filename(headline)
        clean_sqcode = clean_filename(sqcode)
        filename = f"Sample_{clean_market_name}_{clean_sqcode}_Skyquest_2025_V1.pptx"
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)

        logger.info(f"[{request_id}] Saving presentation to: {filepath}")
        prs.save(filepath)

        ppt_elapsed = time.time() - ppt_start_time
        total_elapsed = time.time() - start_time

        timing_logger.info(
            f"[{request_id}] PPT processing completed in {ppt_elapsed:.2f}s"
        )
        timing_logger.info(
            f"[{request_id}] Total request completed in {total_elapsed:.2f}s"
        )

        logger.info(f"[{request_id}] PPT generation completed successfully: {filename}")

        return jsonify(
            {
                "success": True,
                "filename": filename,
                "message": "PowerPoint generated successfully",
            }
        )

    except Exception as e:
        elapsed = time.time() - start_time
        logger.error(
            f"[{request_id}] PPT generation failed after {elapsed:.2f}s: {str(e)}"
        )
        traceback.print_exc()
        return (
            jsonify({"error": "Failed to generate PowerPoint", "message": str(e)}),
            500,
        )


@app.route('/gii/download-gii/<path:filename>')
def download_file(filename):
    download_start_time = time.time()
    logger.info(f"Download request for file: {filename}")

    try:
        import urllib.parse

        decoded_filename = urllib.parse.unquote(filename)

        if (
            ".." in decoded_filename
            or "/" in decoded_filename
            or "\\" in decoded_filename
        ):
            logger.warning(f"Invalid filename attempted: {decoded_filename}")
            return jsonify({"error": "Invalid filename"}), 400

        filepath = os.path.join(app.config["UPLOAD_FOLDER"], decoded_filename)

        if not os.path.exists(filepath):
            logger.warning(f"File not found: {filepath}")
            try:
                available_files = os.listdir(app.config["UPLOAD_FOLDER"])
                logger.info(f"Available files: {available_files}")
            except:
                logger.error("Could not list available files")
            return jsonify({"error": "File not found"}), 404

        def remove_file_after_send(filepath):
            def remove_file(response):
                try:
                    if os.path.exists(filepath):
                        os.remove(filepath)
                        logger.info(f"Temporary file deleted: {filepath}")
                except Exception as e:
                    logger.error(f"Error deleting file: {e}")
                return response

            return remove_file

        response = send_file(
            filepath, as_attachment=True, download_name=decoded_filename
        )

        @response.call_on_close
        def delete_file():
            try:
                if os.path.exists(filepath):
                    os.remove(filepath)
                    logger.info(f"Temporary file deleted after download: {filepath}")
            except Exception as e:
                logger.error(f"Error deleting file after download: {e}")

        elapsed = time.time() - download_start_time
        timing_logger.info(
            f"File download completed in {elapsed:.2f}s: {decoded_filename}"
        )

        return response

    except Exception as e:
        elapsed = time.time() - download_start_time
        logger.error(f"Download failed after {elapsed:.2f}s: {e}")
        return jsonify({"error": "File download failed", "details": str(e)}), 500


if __name__ == "__main__":
    if not os.path.exists(".env"):
        with open(".env", "w") as f:
            f.write("OPENAI_API_KEY=your_openai_api_key_here\n")
            f.write("GEMINI_API_KEY=your_gemini_api_key_here\n")
        logger.info("Created .env file. Please add your API keys.")

    if not os.path.exists("templates/index.html"):
        os.makedirs("templates", exist_ok=True)
        logger.warning(
            "Please save the HTML content from the artifact to templates/index.html"
        )

    logger.info("Starting Flask application")
    app.run(host="0.0.0.0", debug=True, port=5000)